# -*- coding: utf-8 -*-
"""Completa la plantilla IIPSi con el póster BAG CTRL vs pacientes (COMA / COMA+ReDLaT)."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
TEMPLATE = HERE / "Plantilla Póster - 3ra Jornada del IIPSi abre sus puertas.pptx"
OUT = HERE / "Poster_BAG_COMA_IIPSi.pptx"

FIGS_COMA = HERE / "figs_ctrl_pacientes"
FIGS_RL = HERE / "figs_ctrl_pacientes_RedLat"

BURGUNDY = RGBColor(0x5C, 0x1A, 0x2A)
INK = RGBColor(0x2A, 0x22, 0x24)
INK2 = RGBColor(0x4A, 0x3F, 0x42)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xFB, 0xF6, 0xEF)
HAIR = RGBColor(0xD6, 0xCB, 0xC0)

FONT = "Calibri"


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def card(slide, x, y, w, h, fill=CARD):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    solid(sh, fill)
    sh.adjustments[0] = 0.04
    sh.line.color.rgb = HAIR
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def accent_bar(slide, x, y, w, h=Inches(0.08), color=BURGUNDY):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    solid(sh, color)
    sh.shadow.inherit = False
    return sh


def txt(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    spacing=1.05,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = ln
        f = r.font
        f.size = Pt(size)
        f.name = FONT
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def section_title(slide, x, y, w, title: str):
    accent_bar(slide, x, y, Inches(0.55), Inches(0.1))
    txt(
        slide,
        x + Inches(0.7),
        y - Inches(0.05),
        w - Inches(0.7),
        Inches(0.45),
        title.upper(),
        size=22,
        color=BURGUNDY,
        bold=True,
    )


def add_pic(slide, path: Path, x, y, w, max_h=None):
    """Inserta imagen; si max_h, ajusta ancho para no superar esa altura."""
    if not path.exists():
        txt(slide, x, y, w, Inches(0.4), f"[falta {path.name}]", size=14, color=INK2, italic=True)
        return None
    from PIL import Image

    with Image.open(path) as im:
        pw, ph = im.size
    aspect = ph / float(pw)
    height = w * aspect
    if max_h is not None and height > max_h:
        height = max_h
        w = height / aspect
    return slide.shapes.add_picture(str(path), x, y, width=w, height=height)


def clear_empty_title_box(slide):
    for sh in list(slide.shapes):
        if sh.has_text_frame and not sh.text_frame.text.strip():
            sp = sh._element
            sp.getparent().remove(sp)


def build():
    if not TEMPLATE.exists():
        raise FileNotFoundError(TEMPLATE)

    shutil.copy2(TEMPLATE, OUT)
    prs = Presentation(str(OUT))
    slide = prs.slides[0]
    clear_empty_title_box(slide)

    # Título
    txt(
        slide,
        Inches(2.0),
        Inches(8.15),
        Inches(29.0),
        Inches(1.3),
        "Brecha de edad cerebral (BAG) en pacientes post-coma:\n"
        "conectividad funcional y topología de redes",
        size=34,
        color=BURGUNDY,
        bold=True,
        spacing=1.0,
    )
    txt(
        slide,
        Inches(2.0),
        Inches(9.55),
        Inches(29.0),
        Inches(0.45),
        "Ignacio Bosch  ·  TP Final — Neurociencias Computacionales  ·  IIPSi / Facultad de Psicología — UNC",
        size=17,
        color=INK2,
        italic=True,
    )

    # Introducción (texto del usuario, estructurado)
    section_title(slide, Inches(2.0), Inches(10.2), Inches(14.5), "Introducción")
    intro = (
        "La edad cerebral no siempre coincide con la edad cronológica. El estilo de vida, "
        "las enfermedades y cualquier evento neurológico puede acelerar o retrasar el "
        "envejecimiento del cerebro con respecto al paso del tiempo real.\n\n"
        "Esta diferencia se conoce como brecha de edad cerebral (BAG, por sus siglas en inglés), "
        "y se obtiene comparando la edad cronológica de una persona con la edad predicha por un "
        "modelo de inteligencia artificial, a partir de patrones de conectividad cerebral.\n\n"
        "En este trabajo aplicamos este enfoque a pacientes que atravesaron un coma por anoxia "
        "o traumatismo. A partir de datos de resonancia magnética funcional, estimamos la edad "
        "cerebral de cada paciente y calculamos su BAG, buscando caracterizar cómo el tipo de "
        "daño cerebral se refleja en el envejecimiento estimado del cerebro."
    )
    txt(slide, Inches(2.0), Inches(10.75), Inches(14.6), Inches(6.6), intro, size=16.5, color=INK, spacing=1.1)

    # Objetivo
    section_title(slide, Inches(17.2), Inches(10.2), Inches(14.0), "Objetivo")
    objetivo = (
        "Comparar controles sanos frente a pacientes post-coma (anoxia + traumatismo unidos "
        "como un solo grupo), estimando BAG con un modelo entrenado sobre conectividad "
        "funcional y topología de la red cerebral.\n\n"
        "Evaluamos dos diseños:\n"
        "• COMA: controles y pacientes solo del dataset de coma.\n"
        "• COMA+ReDLaT: mismos pacientes, reforzando el grupo control con controles "
        "cognitivamente sanos (CN) de ReDLaT."
    )
    txt(slide, Inches(17.2), Inches(10.75), Inches(14.0), Inches(3.6), objetivo, size=16.5, color=INK, spacing=1.1)

    # Métodos
    section_title(slide, Inches(17.2), Inches(14.6), Inches(14.0), "Métodos")
    metodos = (
        "Como predictores de los modelos de regresión utilizamos tanto la conectividad "
        "funcional como métricas de organización de la red cerebral (topología): eficiencia "
        "local, eficiencia global, agrupamiento y grado promedio de conexión.\n\n"
        "Pipeline:\n"
        "1) FC AAL-116 + Fisher-z.\n"
        "2) Topología de red (umbral 0.20).\n"
        "3) β-VAE → embedding latente Z.\n"
        "4) Ridge (Z + topología + sexo → edad), evaluado por leave-one-out.\n"
        "5) BAG = edad predicha − edad cronológica.\n\n"
        "Se considera que el BAG reporta cambios en la estructura del cerebro y puede "
        "funcionar como un indicador general de su salud, con potencial para detectar "
        "alteraciones de forma temprana y ayudar a definir el diagnóstico, el pronóstico "
        "y el tratamiento más adecuado."
    )
    txt(slide, Inches(17.2), Inches(15.15), Inches(14.0), Inches(5.8), metodos, size=15.5, color=INK, spacing=1.08)

    # Resultados
    section_title(slide, Inches(2.0), Inches(17.7), Inches(29.0), "Resultados")

    card(slide, Inches(2.0), Inches(18.3), Inches(14.5), Inches(2.4))
    card(slide, Inches(17.2), Inches(18.3), Inches(14.0), Inches(2.4))

    txt(slide, Inches(2.2), Inches(18.45), Inches(14.0), Inches(0.35),
        "Diseño A — solo COMA", size=19, color=BURGUNDY, bold=True)
    txt(
        slide,
        Inches(2.2),
        Inches(18.9),
        Inches(14.0),
        Inches(1.6),
        "n = 42  (CTRL 19 · pacientes 23)\n"
        "MAE LOO = 16,9 años  (CTRL 16,4 · pacientes 17,2)\n"
        "BAG medio ≈ 0 en ambos grupos (poca separación con n chico)",
        size=15.5,
        color=INK,
        spacing=1.12,
    )

    txt(slide, Inches(17.4), Inches(18.45), Inches(13.6), Inches(0.35),
        "Diseño B — COMA + controles ReDLaT", size=19, color=BURGUNDY, bold=True)
    txt(
        slide,
        Inches(17.4),
        Inches(18.9),
        Inches(13.6),
        Inches(1.6),
        "n = 597  (CTRL 574 = 19 COMA + 555 CN · pacientes 23)\n"
        "MAE LOO = 8,5 años  (CTRL 8,2 · pacientes 16,0)\n"
        "BAG pacientes ≈ +4,3 años  vs  CTRL ≈ 0",
        size=15.5,
        color=INK,
        spacing=1.12,
    )

    # Figuras fila 1: composición
    txt(slide, Inches(2.0), Inches(20.95), Inches(14.5), Inches(0.32),
        "COMA — composición de la cohorte", size=15, color=INK2, bold=True)
    txt(slide, Inches(17.2), Inches(20.95), Inches(14.0), Inches(0.32),
        "COMA+ReDLaT — composición (control reforzado)", size=15, color=INK2, bold=True)
    add_pic(slide, FIGS_COMA / "01_composicion.png", Inches(2.0), Inches(21.35), Inches(14.4), max_h=Inches(4.5))
    add_pic(slide, FIGS_RL / "01_composicion.png", Inches(17.2), Inches(21.35), Inches(14.0), max_h=Inches(4.5))

    # Figuras fila 2: BAG
    txt(slide, Inches(2.0), Inches(26.0), Inches(14.5), Inches(0.32),
        "COMA — predicción de edad, BAG y error (LOO)", size=15, color=INK2, bold=True)
    txt(slide, Inches(17.2), Inches(26.0), Inches(14.0), Inches(0.32),
        "COMA+ReDLaT — predicción de edad, BAG y error (LOO)", size=15, color=INK2, bold=True)
    add_pic(slide, FIGS_COMA / "04_bag_mae.png", Inches(2.0), Inches(26.4), Inches(14.4), max_h=Inches(4.3))
    add_pic(slide, FIGS_RL / "04_bag_mae.png", Inches(17.2), Inches(26.4), Inches(14.0), max_h=Inches(4.3))

    # Figuras fila 3: topología
    txt(slide, Inches(2.0), Inches(31.0), Inches(14.5), Inches(0.32),
        "COMA — métricas topológicas", size=15, color=INK2, bold=True)
    txt(slide, Inches(17.2), Inches(31.0), Inches(14.0), Inches(0.32),
        "COMA+ReDLaT — métricas topológicas", size=15, color=INK2, bold=True)
    add_pic(slide, FIGS_COMA / "06_topo.png", Inches(2.0), Inches(31.4), Inches(14.4), max_h=Inches(5.5))
    add_pic(slide, FIGS_RL / "06_topo.png", Inches(17.2), Inches(31.4), Inches(14.0), max_h=Inches(5.5))

    # Conclusiones (antes del quote de la plantilla ~40.3")
    section_title(slide, Inches(2.0), Inches(37.2), Inches(29.0), "Conclusiones")
    concl = (
        "Con solo datos de coma, el modelo aprende con poca muestra y el MAE LOO es alto (~17 años): "
        "la señal de BAG apenas separa controles de pacientes. Al reforzar el grupo control con "
        "controles ReDLaT, la predicción de edad mejora de forma clara (MAE global ~8,5 años) y "
        "aparece una diferencia más interpretable: los pacientes post-coma muestran un BAG medio "
        "positivo (~+4 años), compatible con un envejecimiento cerebral acelerado respecto del "
        "control. El BAG, basado en conectividad funcional y topología de red, se perfila como "
        "marcador complementario para caracterizar la salud cerebral tras un evento neurológico grave."
    )
    txt(slide, Inches(2.0), Inches(37.75), Inches(29.0), Inches(1.7), concl, size=16.5, color=INK, spacing=1.1)

    txt(
        slide,
        Inches(2.0),
        Inches(39.5),
        Inches(29.0),
        Inches(0.35),
        "Datos: cohorte inflamación/coma (CTRL, ANOX, TRAU) y CN ReDLaT (data-iipsi). "
        "Modelos entrenados dentro de cada diseño (sin transferir pesos entre estudios).",
        size=12.5,
        color=INK2,
        italic=True,
    )

    prs.save(str(OUT))
    print("guardado:", OUT)


if __name__ == "__main__":
    build()

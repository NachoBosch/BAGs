# -*- coding: utf-8 -*-
"""Genera la presentacion del TP Final (NeuroComp) en estilo editorial sobrio."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

FIG = "figs"

# ---------- Paleta ----------
PAPER   = RGBColor(0xF5, 0xF3, 0xEE)   # blanco calido
INK     = RGBColor(0x20, 0x26, 0x2B)   # carbon
INK2    = RGBColor(0x4A, 0x52, 0x59)   # gris texto
ACCENT  = RGBColor(0x2E, 0x6F, 0x6B)   # teal profundo
ACCENT2 = RGBColor(0xC2, 0x5B, 0x3F)   # terracota
HAIR    = RGBColor(0xD8, 0xD3, 0xC8)   # lineas finas
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CN_C    = RGBColor(0x2E, 0x6F, 0x6B)   # CN teal
AD_C    = RGBColor(0xC2, 0x5B, 0x3F)   # AD terracota
FTD_C   = RGBColor(0xD8, 0x9A, 0x3E)   # FTD ambar
HILITE  = RGBColor(0xE9, 0xE3, 0xD6)   # fondo fila destacada
HEADBG  = INK

FONT = "Segoe UI"
FONT_L = "Segoe UI Light"
FONT_SB = "Segoe UI Semibold"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color=PAPER):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    solid(r, color)
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def txt(slide, x, y, w, h, text, size=18, color=INK, font=FONT, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.name = font; f.bold = bold
        f.italic = italic; f.color.rgb = color
    return tb


def line(slide, x, y, w, h, color=ACCENT):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    solid(r, color); r.shadow.inherit = False
    return r


def header(slide, kicker, title, num):
    bg(slide)
    line(slide, Inches(0.7), Inches(0.62), Inches(0.42), Inches(0.07), ACCENT)
    txt(slide, Inches(1.22), Inches(0.5), Inches(8), Inches(0.3),
        kicker.upper(), size=11.5, color=ACCENT, font=FONT_SB)
    txt(slide, Inches(0.7), Inches(0.78), Inches(11.5), Inches(0.7),
        title, size=27, color=INK, font=FONT_SB)
    line(slide, Inches(0.7), Inches(1.5), Inches(11.93), Pt(1), HAIR)
    # footer
    txt(slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
        "BAGs, topologia y criticalidad  ·  AD / FTD / CN", size=9, color=INK2, font=FONT)
    txt(slide, Inches(11.8), Inches(7.05), Inches(0.83), Inches(0.3),
        str(num), size=9, color=INK2, font=FONT_SB, align=PP_ALIGN.RIGHT)


def add_img(slide, path, x, y, w=None, h=None, maxw=None, maxh=None):
    im = Image.open(path); iw, ih = im.size; ar = iw / ih
    if w is None and h is None:
        if maxw and maxh:
            if maxw / maxh > ar:
                h = maxh; w = Emu(int(maxh * ar))
            else:
                w = maxw; h = Emu(int(maxw / ar))
    if w is not None and h is None:
        h = Emu(int(w / ar))
    if h is not None and w is None:
        w = Emu(int(h * ar))
    return slide.shapes.add_picture(path, x, y, w, h)


def table(slide, x, y, w, rows, colw=None, head=True, fontsize=13,
          hilite_rows=(), align_first_left=True, rowh=0.42, headcolor=HEADBG):
    nrows = len(rows); ncols = len(rows[0])
    h = Inches(rowh * nrows)
    gtbl = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    tbl = gtbl.table
    # kill default style banding
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    if colw:
        total = sum(colw)
        for i, cw in enumerate(colw):
            tbl.columns[i].width = Emu(int(int(w) * cw / total))
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(rowh)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_head = head and ri == 0
            if is_head:
                cell.fill.solid(); cell.fill.fore_color.rgb = headcolor
            elif ri in hilite_rows:
                cell.fill.solid(); cell.fill.fore_color.rgb = HILITE
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else PAPER
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (ci == 0 and align_first_left) else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            f = r.font; f.size = Pt(fontsize); f.name = FONT
            if is_head:
                f.color.rgb = PAPER; f.bold = True; f.size = Pt(fontsize - 0.5)
            else:
                f.color.rgb = INK; f.bold = ri in hilite_rows
    # remove borders, add thin bottom hairlines
    return tbl


def chip(slide, x, y, label, color):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.16), Inches(0.16))
    solid(c, color)
    txt(slide, x + Inches(0.22), y - Inches(0.04), Inches(1.4), Inches(0.3),
        label, size=12, color=INK2, font=FONT_SB)


def legend(slide, x, y):
    chip(slide, x, y, "CN", CN_C)
    chip(slide, x + Inches(1.05), y, "AD", AD_C)
    chip(slide, x + Inches(2.1), y, "FTD", FTD_C)


def note(slide, x, y, w, text, size=11):
    txt(slide, x, y, w, Inches(0.6), text, size=size, color=INK2, font=FONT, italic=True, spacing=1.05)


# ============================================================ SLIDE 1 — Portada
s = prs.slides.add_slide(BLANK)
bg(s, INK)
# accent block
line(s, 0, 0, Inches(0.28), SH, ACCENT)
txt(s, Inches(1.0), Inches(1.05), Inches(11), Inches(0.4),
    "NEUROCIENCIAS COMPUTACIONALES · TP FINAL · POSGRADO", size=13, color=ACCENT, font=FONT_SB)
txt(s, Inches(0.95), Inches(1.9), Inches(11.4), Inches(2.4),
    "BAGs, topologia y criticalidad\nen Alzheimer y Demencia frontotemporal",
    size=42, color=PAPER, font=FONT_L, spacing=1.0)
line(s, Inches(1.0), Inches(4.25), Inches(3.2), Pt(1.5), HAIR)
txt(s, Inches(1.0), Inches(4.5), Inches(11), Inches(1.6),
    "Cohorte latinoamericana ReDLaT–BrainLat  ·  n = 1245  (CN / AD / FTD)\n"
    "fMRI reposo (AAL-116) + T1w (VBM) + β-VAE + teoria de grafos + criticalidad espectral",
    size=15.5, color=RGBColor(0xC9,0xCE,0xCF), font=FONT, spacing=1.25)
txt(s, Inches(1.0), Inches(6.35), Inches(11), Inches(0.6),
    "Ignacio Manuel Bosch   ·   nachobosch5@gmail.com   ·   Junio 2026",
    size=13, color=PAPER, font=FONT_SB)

# ============================================================ SLIDE 2 — Contexto
s = prs.slides.add_slide(BLANK)
header(s, "Contexto y objetivos", "Que problema y que extendemos", 2)
# left: problem
txt(s, Inches(0.7), Inches(1.85), Inches(5.7), Inches(0.4), "EL PROBLEMA", size=12, color=ACCENT, font=FONT_SB)
for i, (t) in enumerate([
    "AD y FTD: carga creciente en paises de ingresos bajos/medios (>70% casos futuros).",
    "America Latina: sub-representada en neuroimagen a gran escala.",
    "Brain Age Gap (BAG) = edad cerebral estimada − edad cronologica.",
    "BAG+ se asocia a deterioro cognitivo y menor reserva."]):
    txt(s, Inches(0.7), Inches(2.3 + i*0.62), Inches(5.7), Inches(0.6),
        "—  " + t, size=13.5, color=INK2, font=FONT, spacing=1.05)
# right: 4 extensiones cards
txt(s, Inches(6.95), Inches(1.85), Inches(5.7), Inches(0.4),
    "ESTE TRABAJO EXTIENDE LA TESIS DE REFERENCIA", size=12, color=ACCENT2, font=FONT_SB)
cards = [
    ("01", "Replica", "Pipeline multimodal β-VAE con Fisher-z; Ridge, SVR, XGB, MLP."),
    ("02", "Topologia (TOPO)", "6 metricas de grafos al vector multimodal → modelos hibridos."),
    ("03", "BAG ↔ topologia", "Comparacion BAG+ vs BAG− y correlacion con MMSE."),
    ("04", "Mundo pequeno + criticalidad", "σ/ω, λ₁ espectral, avalanchas, percolacion."),
]
cy = Inches(2.32)
for code, ti, de in cards:
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.95), cy, Inches(5.68), Inches(0.93))
    solid(card, WHITE); card.line.color.rgb = HAIR; card.line.width = Pt(0.75)
    card.shadow.inherit = False
    txt(s, Inches(7.12), cy + Inches(0.13), Inches(0.8), Inches(0.6), code, size=22, color=ACCENT, font=FONT_L)
    txt(s, Inches(7.95), cy + Inches(0.1), Inches(4.6), Inches(0.35), ti, size=14.5, color=INK, font=FONT_SB)
    txt(s, Inches(7.95), cy + Inches(0.46), Inches(4.6), Inches(0.4), de, size=11, color=INK2, font=FONT)
    cy += Inches(1.04)

# ============================================================ SLIDE 3 — Metodos / pipeline
s = prs.slides.add_slide(BLANK)
header(s, "Metodos", "Pipeline multimodal de edad cerebral", 3)
steps = [
    ("FC  116×116", "Pearson → Fisher-z\n6670 conexiones"),
    ("β-VAE", "embeddings\nZ ∈ R⁶⁴"),
    ("Fusion", "Z + T1w(VBM) + sexo\n+ TOPO + educ. + sitio"),
    ("Regresores", "Ridge · SVR · XGB · MLP\nHP post-Optuna"),
    ("BAG", "entrenado en CN\ncorreccion LOSO por sitio"),
]
n = len(steps); x0 = Inches(0.7); gap = Inches(0.25)
cardw = Emu(int((int(SW) - int(Inches(1.4)) - int(gap)*(n-1)) / n))
for i, (ti, de) in enumerate(steps):
    cx = Emu(int(x0) + i*(int(cardw)+int(gap)))
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(2.55), cardw, Inches(1.7))
    card.adjustments[0] = 0.08
    solid(card, WHITE); card.line.color.rgb = HAIR; card.line.width = Pt(1); card.shadow.inherit = False
    line(s, cx, Inches(2.55), cardw, Inches(0.09), ACCENT)
    txt(s, cx + Inches(0.12), Inches(2.78), Emu(int(cardw)-int(Inches(0.24))), Inches(0.5),
        ti, size=14.5, color=INK, font=FONT_SB, align=PP_ALIGN.CENTER)
    txt(s, cx + Inches(0.12), Inches(3.35), Emu(int(cardw)-int(Inches(0.24))), Inches(0.8),
        de, size=11, color=INK2, font=FONT, align=PP_ALIGN.CENTER, spacing=1.05)
    if i < n-1:
        txt(s, Emu(int(cx)+int(cardw)-int(Inches(0.02))), Inches(3.05), Inches(0.3), Inches(0.4),
            "›", size=22, color=ACCENT, font=FONT_SB, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4), "DECISIONES CLAVE", size=12, color=ACCENT, font=FONT_SB)
for i, t in enumerate([
    "Umbral fijo Fisher-z = 0.20 para grafos · umbral proporcional 10% para criticalidad.",
    "Holdout 90/10 estratificado (sexo × diagnostico); efectos de sitio fuertes (Kruskal-Wallis p<0.001) → sitio como covariable.",
    "Se reutilizan HP post-Optuna de la referencia de forma uniforme (sin re-busqueda por configuracion)."]):
    txt(s, Inches(0.7), Inches(5.15 + i*0.5), Inches(12), Inches(0.5), "—  " + t, size=13, color=INK2, font=FONT)

# ============================================================ SLIDE 4 — Muestra
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · Muestra", "Cohorte y demografia", 4)
rows = [
    ["", "CN", "AD", "FTD", "Total"],
    ["N", "555", "468", "304", "1327"],
    ["Edad (media ± DE)", "63.3 ± 11.4", "70.0 ± 8.5", "65.3 ± 8.3", "65.9 ± 10.5"],
    ["Sitios con datos", "9", "9", "9", "9"],
]
table(s, Inches(0.7), Inches(1.95), Inches(6.4), rows, colw=[2.3,1,1,1,1], fontsize=13.5, rowh=0.6)
legend(s, Inches(0.7), Inches(4.6))
note(s, Inches(0.7), Inches(5.0), Inches(6.4),
     "Diferencias de edad entre grupos (Kruskal-Wallis, p<0.001); AD el grupo mas anoso.\n"
     "9 sitios; efectos de sitio sobre edad (H=312.4) y educacion (H=287.1), p<0.001.")
add_img(s, FIG + "/01_age_distributions.png", Inches(7.4), Inches(2.3), maxw=Inches(5.5), maxh=Inches(3.6))
txt(s, Inches(7.4), Inches(2.0), Inches(5.5), Inches(0.3), "Distribucion de edad por grupo", size=12, color=INK2, font=FONT_SB)

# ============================================================ SLIDE 5 — Metricas topologicas
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · Topologia", "Seis metricas de grafos por diagnostico", 5)
rows = [
    ["Metrica", "CN", "AD", "FTD"],
    ["E local", "0.861 ± .057", "0.849 ± .049", "0.858 ± .058"],
    ["E global", "0.763", "0.744", "0.755"],
    ["Clustering", "0.726", "0.694", "0.714"],
    ["E local FP", "0.861", "0.847", "0.861"],
    ["Grado FP", "0.599 ± .19", "0.548 ± .17", "0.575 ± .19"],
]
table(s, Inches(0.7), Inches(1.95), Inches(5.6), rows, colw=[1.9,1,1,1], fontsize=12.5, rowh=0.5)
note(s, Inches(0.7), Inches(5.0), Inches(5.6),
     "AD se desplaza a valores menores en todas las metricas (eficiencias, clustering, hubs FP).\nUmbral Fisher-z = 0.20.")
add_img(s, FIG + "/02_graph_metrics_by_diagnosis.png", Inches(6.6), Inches(1.95), maxw=Inches(6.2), maxh=Inches(4.6))

# ============================================================ SLIDE 6 — Prediccion: referencia vs replica + ablacion
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · Prediccion", "Replica de la referencia y ablacion", 6)
txt(s, Inches(0.7), Inches(1.75), Inches(6), Inches(0.3), "REFERENCIA vs REPLICA", size=12, color=ACCENT, font=FONT_SB)
rows = [
    ["Configuracion", "n", "MAE", "R²", "r"],
    ["Ref.: Z+T1w+educ+sitio (XGB+Optuna)", "91", "5.17", ".523", ".725"],
    ["Este trabajo: Z+T1w+educ+sitio", "91", "5.54", ".473", ".693"],
    ["Este trabajo: Z+T1w+sexo", "125", "5.88", ".407", ".640"],
]
table(s, Inches(0.7), Inches(2.1), Inches(6.0), rows, colw=[3.2,0.7,0.9,0.9,0.8], fontsize=11, rowh=0.52, hilite_rows=(1,))
txt(s, Inches(0.7), Inches(4.55), Inches(6), Inches(0.3), "ABLACION DE CARACTERISTICAS", size=12, color=ACCENT, font=FONT_SB)
rows = [
    ["Config.", "Dims", "MAE", "R²"],
    ["T1w solo", "116", "5.85", ".385"],
    ["Z solo", "64", "7.11", ".102"],
    ["Z + T1w", "180", "5.79", ".417"],
    ["Z + T1w + sexo + dx", "182", "5.74", ".419"],
]
table(s, Inches(0.7), Inches(4.9), Inches(6.0), rows, colw=[2.4,0.8,0.8,0.8], fontsize=11, rowh=0.36)
# right column - takeaways
txt(s, Inches(7.2), Inches(1.85), Inches(5.4), Inches(0.4), "LECTURA", size=12, color=ACCENT2, font=FONT_SB)
for i, (big, t) in enumerate([
    ("≈5.5–5.9", "MAE de la replica en holdout, mismo orden que la referencia (5.17)."),
    ("5.79", "Z + T1w confirma complementariedad conectividad ⊕ morfometria."),
    ("7.11", "Z solo es el peor predictor; T1w solo ya da 5.85."),
]):
    yy = Inches(2.35 + i*1.35)
    txt(s, Inches(7.2), yy, Inches(2.4), Inches(0.7), big, size=34, color=INK, font=FONT_L)
    txt(s, Inches(9.7), yy + Inches(0.08), Inches(2.95), Inches(1.1), t, size=12.5, color=INK2, font=FONT, spacing=1.1)

# ============================================================ SLIDE 7 — TOPO + Modelo hibrido
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · Prediccion", "Topologia hibrida → mejor modelo del estudio", 7)
txt(s, Inches(0.7), Inches(1.75), Inches(6), Inches(0.3), "APORTE DE TOPO (sobre Z+T1w+sexo)", size=12, color=ACCENT, font=FONT_SB)
rows = [
    ["Configuracion", "MAE", "R²", "r"],
    ["TOPO sola", "8.21", "−0.22", ".098"],
    ["Z + T1w + sexo  (base)", "5.88", ".407", ".640"],
    ["Z + T1w + sexo + TOPO", "5.76", ".412", ".646"],
]
table(s, Inches(0.7), Inches(2.1), Inches(5.9), rows, colw=[3,0.9,0.9,0.8], fontsize=11.5, rowh=0.5)
txt(s, Inches(0.7), Inches(4.35), Inches(6.3), Inches(0.4),
    "MODELO HIBRIDO COMPLETO   Z+T1w+sexo+TOPO+educ+sitio  (n=91, 197 feats)",
    size=12, color=ACCENT2, font=FONT_SB)
rows = [
    ["Regresor", "MAE", "R²", "r"],
    ["Ridge  (mejor global)", "5.06", ".536", ".736"],
    ["SVR", "5.13", ".547", ".742"],
    ["XGBoost (HP unicos)", "5.44", ".488", ".705"],
    ["Referencia (sin TOPO)", "5.17", ".523", ".725"],
]
table(s, Inches(0.7), Inches(4.75), Inches(5.9), rows, colw=[3,0.9,0.9,0.8], fontsize=11.5, rowh=0.42, hilite_rows=(1,))
# big number highlight
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.0))
box.adjustments[0] = 0.06
solid(box, INK); box.shadow.inherit = False
txt(s, Inches(7.3), Inches(2.18), Inches(5.1), Inches(0.4), "RIDGE HIBRIDO · MAE EN TEST", size=12, color=ACCENT, font=FONT_SB)
txt(s, Inches(7.25), Inches(2.5), Inches(3.2), Inches(1.2), "5.06", size=66, color=PAPER, font=FONT_L)
txt(s, Inches(10.4), Inches(2.75), Inches(2.2), Inches(1), "anos\nR² = 0.54", size=16, color=RGBColor(0xC9,0xCE,0xCF), font=FONT, spacing=1.1)
txt(s, Inches(7.3), Inches(3.62), Inches(5.1), Inches(0.4),
    "▼ 0.11 anos bajo la mejor cifra de la referencia (5.17, sin TOPO)", size=12.5, color=FTD_C, font=FONT_SB)
add_img(s, FIG + "/06_test_scatter_best.png", Inches(7.5), Inches(4.25), maxw=Inches(4.8), maxh=Inches(2.7))
note(s, Inches(0.7), Inches(6.65), Inches(11), "MLP (3 capas densas): MAE 8.3, R²<0 → sobreajuste; se prefieren modelos lineales/kernel en este regimen.", size=10.5)

# ============================================================ SLIDE 8 — BAG por diagnostico + BAG+/-
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · BAG", "Brain Age Gap: diagnostico y topologia", 8)
txt(s, Inches(0.7), Inches(1.75), Inches(6), Inches(0.3), "BAG MEDIO POR DIAGNOSTICO (corregido)", size=12, color=ACCENT, font=FONT_SB)
for i, (g, v, c) in enumerate([("CN", "1.70 ± 6.0", CN_C), ("AD", "11.58 ± 6.8", AD_C), ("FTD", "13.93 ± 6.4", FTD_C)]):
    cx = Inches(0.7 + i*2.05)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(2.15), Inches(1.85), Inches(1.25))
    card.adjustments[0]=0.08; solid(card, WHITE); card.line.color.rgb=HAIR; card.line.width=Pt(1); card.shadow.inherit=False
    line(s, cx, Inches(2.15), Inches(1.85), Inches(0.08), c)
    txt(s, cx, Inches(2.3), Inches(1.85), Inches(0.3), g, size=14, color=c, font=FONT_SB, align=PP_ALIGN.CENTER)
    txt(s, cx, Inches(2.62), Inches(1.85), Inches(0.6), v.split(" ")[0], size=30, color=INK, font=FONT_L, align=PP_ALIGN.CENTER)
    txt(s, cx, Inches(3.18), Inches(1.85), Inches(0.3), "anos", size=10.5, color=INK2, font=FONT, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(3.7), Inches(6.3), Inches(0.4), "BAG+ vs BAG−   (Mann-Whitney U, FDR-BH)", size=12, color=ACCENT2, font=FONT_SB)
rows = [
    ["Metrica", "BAG+", "BAG−", "p(FDR)"],
    ["E local", "0.868", "0.856", "0.021 *"],
    ["E global", "0.773", "0.756", "0.039 *"],
    ["Clustering", "0.741", "0.718", "0.021 *"],
    ["E local FP", "0.866", "0.856", "0.039 *"],
    ["Betweenness FP", "0.0048", "0.0051", "0.126"],
    ["Grado FP", "0.612", "0.588", "0.126"],
]
table(s, Inches(0.7), Inches(4.1), Inches(6.3), rows, colw=[2.2,1,1,1.1], fontsize=11.5, rowh=0.36, hilite_rows=(1,2,3,4))
add_img(s, FIG + "/08_topology_BAG_groups.png", Inches(7.3), Inches(2.0), maxw=Inches(5.5), maxh=Inches(4.2))
note(s, Inches(7.3), Inches(6.2), Inches(5.5),
     "Direccion positiva (mas eficiencia en BAG+) coherente con artefacto de umbral fijo: mayor FC → grafos mas densos.", size=10.5)

# ============================================================ SLIDE 9 — MMSE / regresion
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · BAG", "Asociacion con cognicion (MMSE)", 9)
txt(s, Inches(0.7), Inches(1.85), Inches(6), Inches(0.3),
    "CORRELACION PARCIAL SPEARMAN  topologia ↔ MMSE  (control edad+educ, n=916)", size=12, color=ACCENT, font=FONT_SB)
rows = [
    ["Metrica", "r", "p(FDR)"],
    ["Grado FP", "0.104", "0.010 *"],
    ["E global", "0.086", "0.019 *"],
    ["E local", "0.082", "0.019 *"],
]
table(s, Inches(0.7), Inches(2.3), Inches(5.2), rows, colw=[2,1,1.2], fontsize=13, rowh=0.5, hilite_rows=(1,2,3))
note(s, Inches(0.7), Inches(4.7), Inches(5.4),
     "Asociaciones positivas pero debiles (r≈0.08–0.10).\n"
     "Regresion multiple TOPO ~ BAG+covariables: solo grado FP marginal con BAG "
     "(β=−0.0019, p=0.037, FDR=0.22).")
add_img(s, FIG + "/08_topo_mmse_correlations.png", Inches(6.4), Inches(2.4), maxw=Inches(6.3), maxh=Inches(3.6))

# ============================================================ SLIDE 10 — Small world
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · Mundo pequeno", "Small-worldness vs referencias WS y ER", 10)
rows = [
    ["Grupo", "σ WS", "ω WS", "σ ER", "ω ER"],
    ["CN", "1.187", "0.077", "1.467", "0.076"],
    ["AD", "1.221", "0.086", "1.510", "0.085"],
    ["FTD", "1.092", "0.028", "1.305", "0.027"],
]
table(s, Inches(0.7), Inches(2.0), Inches(6.1), rows, colw=[1.2,1,1,1,1], fontsize=12.5, rowh=0.5)
legend(s, Inches(0.7), Inches(4.4))
note(s, Inches(0.7), Inches(4.85), Inches(6.1),
     "σ>1 en TODOS los grupos y ambas referencias → small-world robusto.\n"
     "σ sin diferencias entre grupos (KW p>0.10).  ω SI difiere (KW p≈0.03):\n"
     "FTD con ω menor que CN y AD (post-hoc FDR<0.05) → mas aleatorio.\n"
     "n=30 (10/grupo, estratificado por sitio), N_ref=20.")
add_img(s, FIG + "/08_sw_by_diagnosis.png", Inches(7.0), Inches(2.4), maxw=Inches(5.9), maxh=Inches(3.6))

# ============================================================ SLIDE 11 — Criticalidad lambda1
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · Criticalidad", "Autovalor maximo λ₁ y rango dinamico", 11)
rows = [
    ["Grupo", "λ₁", "λ₁ norm", "n"],
    ["CN", "20.96 ± 3.1", "1.69", "526"],
    ["AD", "20.09 ± 2.6", "1.62", "422"],
    ["FTD", "20.98 ± 3.0", "1.69", "297"],
]
table(s, Inches(0.7), Inches(2.0), Inches(5.5), rows, colw=[1.2,1.4,1.1,0.9], fontsize=12.5, rowh=0.5, hilite_rows=(2,))
txt(s, Inches(0.7), Inches(4.25), Inches(5.6), Inches(0.4), "KRUSKAL-WALLIS   H = 21.35,  p = 2.3·10⁻⁵", size=12.5, color=ACCENT2, font=FONT_SB)
note(s, Inches(0.7), Inches(4.65), Inches(5.6),
     "Post-hoc FDR-BH:\n"
     "•  AD < CN   (p_FDR = 0.00014)\n"
     "•  AD < FTD  (p_FDR = 0.00014)\n"
     "•  CN = FTD  (p_FDR = 0.81)\n"
     "Rango dinamico (matrices promedio): AD el mayor (17.46 dB) pese a menor λ₁ → red promedio menos supercritica.")
add_img(s, FIG + "/08_crit_lambda1.png", Inches(6.5), Inches(2.4), maxw=Inches(6.3), maxh=Inches(3.5))

# ============================================================ SLIDE 12 — Avalanchas + branching
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · Criticalidad", "Avalanchas neuronales y ramificacion", 12)
rows = [
    ["Grupo", "⟨s⟩", "α(s)", "α(T)", "σ_br  [IC95%]"],
    ["CN", "2.97", "2.01", "2.39", "0.663 [.659,.667]"],
    ["AD", "2.70", "2.12", "2.51", "0.630 [.626,.634]"],
    ["FTD", "2.91", "2.02", "2.40", "0.657 [.653,.661]"],
]
table(s, Inches(0.7), Inches(2.0), Inches(6.6), rows, colw=[1,0.9,0.9,0.9,2.3], fontsize=11.5, rowh=0.5)
note(s, Inches(0.7), Inches(4.4), Inches(6.6),
     "Simulacion 50 000 avalanchas, p = 1/λ₁(CN) ≈ 0.049 (fija).\n"
     "α(s) > 2 en los tres grupos (sobre el 1.5 de campo medio).\n"
     "σ_br < 1 → regimen SUBCRITICO bajo esa p. AD: menor ⟨s⟩ y mayor α(s).")
add_img(s, FIG + "/08_crit_avalanchas.png", Inches(7.4), Inches(2.5), maxw=Inches(5.5), maxh=Inches(3.4))

# ============================================================ SLIDE 13 — Percolacion + sensibilidad
s = prs.slides.add_slide(BLANK)
header(s, "Resultados · Criticalidad", "Percolacion y robustez al umbral", 13)
txt(s, Inches(0.7), Inches(1.8), Inches(6), Inches(0.3), "PERCOLACION  (n=30, 10/grupo)", size=12, color=ACCENT, font=FONT_SB)
rows = [
    ["Grupo", "pico χ (%)", "Tc medio (%)"],
    ["CN", "1.8", "2.7"],
    ["AD", "2.8", "2.1"],
    ["FTD", "3.0", "4.2"],
]
table(s, Inches(0.7), Inches(2.2), Inches(4.8), rows, colw=[1.2,1.4,1.4], fontsize=12.5, rowh=0.45)
note(s, Inches(0.7), Inches(4.3), Inches(5.0),
     "Sin diferencias significativas entre grupos en S2_peak, χ_peak ni Tc (KW p>0.24). Alta variabilidad intra-grupo.")
add_img(s, FIG + "/08_crit_percolacion.png", Inches(0.7), Inches(5.0), maxw=Inches(5.6), maxh=Inches(1.9))
txt(s, Inches(6.7), Inches(1.8), Inches(6), Inches(0.3), "SENSIBILIDAD DE λ₁ AL UMBRAL (2–20%)", size=12, color=ACCENT2, font=FONT_SB)
add_img(s, FIG + "/08_crit_sensibilidad.png", Inches(6.9), Inches(2.3), maxw=Inches(5.9), maxh=Inches(3.9))
note(s, Inches(6.7), Inches(6.35), Inches(5.9),
     "AD por debajo de CN y FTD en TODO el rango, sin cruzarse. CN≈FTD en cada punto.", size=10.5)

# ============================================================ SLIDE 14 — Conclusiones
s = prs.slides.add_slide(BLANK)
bg(s, INK)
line(s, 0, 0, Inches(0.28), SH, ACCENT)
txt(s, Inches(0.95), Inches(0.6), Inches(10), Inches(0.4), "CONCLUSIONES", size=13, color=ACCENT, font=FONT_SB)
txt(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.7), "Que aprendimos", size=30, color=PAPER, font=FONT_L)
items = [
    ("5.06", "anos MAE", "Ridge hibrido (Z+T1w+sexo+TOPO+educ+sitio) supera la referencia (5.17). TOPO sola predice mal (~8.2) pero suma integrada."),
    ("FDR<.05", "BAG ↔ topologia", "BAG separa diagnosticos (CN≈1.7 vs AD/FTD≈12). BAG+ con mayor eficiencia/clustering (efecto de umbral fijo a considerar)."),
    ("σ>1", "mundo pequeno", "Small-world robusto en todos los grupos (WS y ER). FTD con ω menor → mas aleatorio."),
    ("AD↓λ₁", "criticalidad", "AD con λ₁ menor que CN y FTD (p_FDR<0.001), estable 2–20%. Regimen subcritico (σ_br<1, α>2). Percolacion sin diferencias."),
]
yy = Inches(2.0)
for big, lab, de in items:
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), yy, Inches(11.5), Inches(1.12))
    solid(card, RGBColor(0x2A,0x31,0x37)); card.line.fill.background(); card.shadow.inherit=False
    line(s, Inches(0.95), yy, Inches(0.06), Inches(1.12), ACCENT)
    txt(s, Inches(1.2), yy + Inches(0.14), Inches(2.3), Inches(0.85), big, size=32, color=PAPER, font=FONT_L)
    txt(s, Inches(3.5), yy + Inches(0.13), Inches(2.6), Inches(0.9), lab, size=15, color=FTD_C, font=FONT_SB, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(6.1), yy + Inches(0.13), Inches(6.1), Inches(0.9), de, size=12, color=RGBColor(0xD2,0xD6,0xD7), font=FONT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)
    yy += Inches(1.22)

prs.save("Presentacion_TPFinal_NeuroComp.pptx")
print("OK -> Presentacion_TPFinal_NeuroComp.pptx  |  slides:", len(prs.slides._sldIdLst))
